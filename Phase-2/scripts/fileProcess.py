
from urllib.parse import unquote, urlparse
from pptx import Presentation
import spacy
import json
import re
from textProcess import TextProcessor
import pdfplumber
# def readPDF(pdf_path):
#     text = ""
#     path=unquote(urlparse(pdf_path).path)
#     pdf_path=path
#     with open(pdf_path, 'rb') as file:
#         readingFile = PyPDF2.PdfReader(file)
#         for pgnum in range(len(readingFile.pages)):
#             page = readingFile.pages[pgnum]
#             extracted_text = page.extract_text()
#             # Clean up extracted text
#             cleaned_text = clean_text(extracted_text)
#             text += cleaned_text + " "  # Append cleaned text with space
    
#     return text.strip()  # Strip leading/trailing whitespaces
def read_PDF(path):
    file_path=unquote(urlparse(path).path)
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text()+"."
    return clean_text(text)

def clean_text( text):
    cleaned_text = re.sub(r'[^\x00-\x7F]+', '', text)
    cleaned_text = re.sub(r'\s+', ' ', cleaned_text)
    cleaned_text = cleaned_text.replace('-\n', '. ')
    return cleaned_text

def read_pptx(uri):
    path=unquote(urlparse(uri).path)
    #print(path)
    prs=Presentation(path)
    texts=[]
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    texts.append(run.text)

    return clean_text(". ".join(texts))


def processFile(uri:str,tp:TextProcessor):
    print("uri is ",uri)
    if uri.endswith(".pptx") or uri.endswith('.ppt'):
        text=read_pptx(uri)
        db,wiki=tp.processText(text)
        return db, wiki
    elif uri.endswith(".pdf"):
        text=read_PDF(uri)
        db,wiki=tp.processText(text)
        return db, wiki
        


if __name__=="__main__":
    #text=read_PDF("file:////Users/hari/Intelligent_Systems_Project/Phase-1/lecture_content/COMP474_6741/slides01.pdf")
    #print(text)
    text=read_PDF("file:////Users/hari/Intelligent_Systems_Project/Phase-2/lecture_content/COMP474_6741/lab2.pdf")
    print(text)
    sp=TextProcessor()
    db,wiki=sp.processText(text)

    print(db)
    print(wiki)

